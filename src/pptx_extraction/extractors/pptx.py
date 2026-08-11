"""PowerPoint OOXML extraction into pptx_extraction domain records."""

from __future__ import annotations

import hashlib
import logging
from collections.abc import Iterable
from datetime import date, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..exceptions import ExtractionError, OptionalDependencyError
from ..models import (
    BoundingBox,
    ChartRecord,
    ChartSeries,
    ExtractionOptions,
    ImageRecord,
    Issue,
    PresentationRecord,
    SlideRecord,
    TableRecord,
    TextBlock,
)
from ..ocr import OCRBackend, create_ocr_backend
from ..security import PackageLimits, sha256_file, validate_package

logger = logging.getLogger(__name__)


class PptxExtractor:
    """Extract one validated OOXML presentation without network access."""

    def __init__(
        self,
        options: ExtractionOptions | None = None,
        *,
        package_limits: PackageLimits | None = None,
        ocr_backend: OCRBackend | None = None,
        tesseract_command: str | None = None,
    ) -> None:
        self.options = options or ExtractionOptions()
        self.package_limits = package_limits or PackageLimits()
        self.ocr = ocr_backend or create_ocr_backend(self.options.ocr_backend, tesseract_command)
        self._ocr_cache: dict[str, str] = {}
        self._asset_cache: dict[str, str] = {}

    def extract(self, source: str | Path, asset_dir: Path | None = None) -> PresentationRecord:
        path = Path(source).expanduser().resolve()
        report = validate_package(path, self.package_limits)
        warnings = [Issue("macro_present", item) for item in report.warnings]
        self._ocr_cache.clear()
        self._asset_cache.clear()
        try:
            presentation = Presentation(str(path))
            slide_width = int(presentation.slide_width or 0)
            slide_height = int(presentation.slide_height or 0)
            if asset_dir and self.options.include_assets:
                asset_dir.mkdir(parents=True, exist_ok=True)
            slides = [
                self._extract_slide(
                    slide,
                    number,
                    slide_width,
                    slide_height,
                    asset_dir if self.options.include_assets else None,
                    warnings,
                )
                for number, slide in enumerate(presentation.slides, start=1)
            ]
            metadata = self._extract_metadata(presentation) if self.options.include_metadata else {}
            if self.options.redact_metadata:
                metadata = self._redact_metadata(metadata)
            return PresentationRecord(
                source_name=path.name,
                source_sha256=sha256_file(path),
                source_size_bytes=path.stat().st_size,
                slide_width_pt=round(slide_width / 12_700, 3),
                slide_height_pt=round(slide_height / 12_700, 3),
                metadata=metadata,
                slides=slides,
                warnings=warnings,
            )
        except OptionalDependencyError:
            raise
        except Exception as exc:
            raise ExtractionError(f"Unable to extract {path.name}: {exc}") from exc

    def _extract_slide(
        self,
        slide: Any,
        number: int,
        slide_width: int,
        slide_height: int,
        asset_dir: Path | None,
        warnings: list[Issue],
    ) -> SlideRecord:
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape and title_shape.text.strip() else None
        record = SlideRecord(
            number=number,
            title=title,
            hidden=slide._element.get("show") == "0",
            layout_name=getattr(slide.slide_layout, "name", None),
            notes=self._extract_notes(slide, warnings, number)
            if self.options.include_notes
            else None,
        )
        flattened = list(self._walk_shapes(slide.shapes))
        row_band = max(int(slide_height * 0.015), 1)
        ordered = sorted(
            flattened,
            key=lambda item: (
                int(getattr(item[1], "top", 0)) // row_band,
                int(getattr(item[1], "left", 0)),
                item[0],
            ),
        )
        for reading_order, (z_order, shape) in enumerate(ordered, start=1):
            bbox = self._bbox(shape, slide_width, slide_height)
            shape_id = int(getattr(shape, "shape_id", 0))
            shape_name = str(getattr(shape, "name", f"shape-{shape_id}"))
            try:
                if getattr(shape, "has_table", False):
                    rows = tuple(
                        tuple(cell.text.strip() for cell in row.cells) for row in shape.table.rows
                    )
                    record.tables.append(
                        TableRecord(rows, reading_order, z_order, shape_id, shape_name, bbox)
                    )
                elif getattr(shape, "has_chart", False):
                    record.charts.append(
                        self._extract_chart(
                            shape, reading_order, z_order, shape_id, shape_name, bbox
                        )
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    record.images.append(
                        self._extract_image(
                            shape,
                            reading_order,
                            z_order,
                            shape_id,
                            shape_name,
                            bbox,
                            asset_dir,
                            number,
                            warnings,
                        )
                    )
                elif getattr(shape, "has_text_frame", False):
                    is_title = bool(
                        title_shape is not None and shape._element is title_shape._element
                    )
                    record.text_blocks.extend(
                        self._extract_text_blocks(
                            shape,
                            "title" if is_title else "body",
                            reading_order,
                            z_order,
                            shape_id,
                            shape_name,
                            bbox,
                        )
                    )
                elif self._is_material_unsupported_shape(shape):
                    warnings.append(
                        Issue(
                            "unsupported_shape",
                            f"Shape type {shape.shape_type} was not extracted.",
                            slide_number=number,
                            shape_id=shape_id,
                        )
                    )
            except OptionalDependencyError:
                raise
            except Exception as exc:
                warnings.append(
                    Issue(
                        "shape_extraction_failed",
                        f"{shape_name}: {exc}",
                        slide_number=number,
                        shape_id=shape_id,
                    )
                )
        return record

    def _walk_shapes(self, shapes: Iterable[Any]) -> Iterable[tuple[int, Any]]:
        z_order = 0
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for _, child in self._walk_shapes(shape.shapes):
                    yield z_order, child
                    z_order += 1
            else:
                yield z_order, shape
                z_order += 1

    @staticmethod
    def _bbox(shape: Any, slide_width: int, slide_height: int) -> BoundingBox:
        return BoundingBox.from_emu(
            int(getattr(shape, "left", 0)),
            int(getattr(shape, "top", 0)),
            int(getattr(shape, "width", 0)),
            int(getattr(shape, "height", 0)),
            slide_width,
            slide_height,
        )

    @staticmethod
    def _extract_text_blocks(
        shape: Any,
        kind: str,
        order: int,
        z_order: int,
        shape_id: int,
        shape_name: str,
        bbox: BoundingBox,
    ) -> list[TextBlock]:
        blocks: list[TextBlock] = []
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            links: list[str] = []
            for run in paragraph.runs:
                try:
                    address = run.hyperlink.address
                except (KeyError, ValueError):
                    address = None
                if address and address not in links:
                    links.append(address)
            blocks.append(
                TextBlock(
                    text=text,
                    kind=kind,
                    level=int(getattr(paragraph, "level", 0)),
                    order=order,
                    z_order=z_order,
                    shape_id=shape_id,
                    shape_name=shape_name,
                    bbox=bbox,
                    hyperlinks=tuple(links),
                )
            )
        return blocks

    @staticmethod
    def _extract_chart(
        shape: Any,
        order: int,
        z_order: int,
        shape_id: int,
        shape_name: str,
        bbox: BoundingBox,
    ) -> ChartRecord:
        chart = shape.chart
        title: str | None = None
        if getattr(chart, "has_title", False):
            candidate = chart.chart_title.text_frame.text.strip()
            title = candidate or None
        categories: list[str] = []
        try:
            if chart.plots:
                for category in chart.plots[0].categories:
                    label = getattr(category, "label", category)
                    categories.append(str(label))
        except (AttributeError, TypeError, ValueError):
            categories = []
        series = tuple(
            ChartSeries(
                name=str(item.name or ""),
                values=tuple(PptxExtractor._json_scalar(value) for value in item.values),
            )
            for item in chart.series
        )
        chart_type = getattr(getattr(chart, "chart_type", None), "name", None)
        return ChartRecord(
            chart_type=chart_type or str(chart.chart_type),
            title=title,
            categories=tuple(categories),
            series=series,
            order=order,
            z_order=z_order,
            shape_id=shape_id,
            shape_name=shape_name,
            bbox=bbox,
        )

    def _extract_image(
        self,
        shape: Any,
        order: int,
        z_order: int,
        shape_id: int,
        shape_name: str,
        bbox: BoundingBox,
        asset_dir: Path | None,
        slide_number: int,
        warnings: list[Issue],
    ) -> ImageRecord:
        image = shape.image
        blob = image.blob
        digest = hashlib.sha256(blob).hexdigest()
        extension = (image.ext or "bin").lower().lstrip(".")
        media_extension = "jpeg" if extension in {"jpg", "jpeg"} else extension
        media_type = f"image/{media_extension}"
        asset_path: str | None = None
        if asset_dir is not None:
            asset_path = self._asset_cache.get(digest)
            if asset_path is None:
                filename = f"{digest[:16]}.{extension}"
                destination = asset_dir / filename
                if not destination.exists():
                    destination.write_bytes(blob)
                asset_path = f"assets/{filename}"
                self._asset_cache[digest] = asset_path
        alt_text = self._alt_text(shape)
        if not alt_text:
            warnings.append(
                Issue(
                    "image_missing_alt_text",
                    f"Image {shape_name} has no alternative text.",
                    slide_number=slide_number,
                    shape_id=shape_id,
                )
            )
        ocr_text = self._ocr_cache.get(digest)
        if ocr_text is None:
            ocr_text = self.ocr.recognize(blob, self.options.ocr_language).strip()
            self._ocr_cache[digest] = ocr_text
        return ImageRecord(
            sha256=digest,
            media_type=media_type,
            asset_path=asset_path,
            alt_text=alt_text,
            ocr_text=ocr_text or None,
            order=order,
            z_order=z_order,
            shape_id=shape_id,
            shape_name=shape_name,
            bbox=bbox,
        )

    @staticmethod
    def _alt_text(shape: Any) -> str | None:
        try:
            nodes = shape._element.xpath(".//p:cNvPr")
            if nodes:
                value = nodes[0].get("descr") or nodes[0].get("title")
                return value.strip() if value and value.strip() else None
        except (AttributeError, IndexError, TypeError):
            return None
        return None

    @staticmethod
    def _extract_notes(slide: Any, warnings: list[Issue], slide_number: int) -> str | None:
        try:
            if not slide.has_notes_slide:
                return None
            text_frame = slide.notes_slide.notes_text_frame
            text = text_frame.text.strip() if text_frame is not None else ""
            return text or None
        except Exception as exc:
            warnings.append(
                Issue(
                    "notes_extraction_failed",
                    str(exc),
                    slide_number=slide_number,
                )
            )
            return None

    @staticmethod
    def _extract_metadata(presentation: Any) -> dict[str, Any]:
        props = presentation.core_properties
        values = {
            "title": props.title,
            "author": props.author,
            "subject": props.subject,
            "keywords": props.keywords,
            "comments": props.comments,
            "last_modified_by": props.last_modified_by,
            "created": props.created,
            "modified": props.modified,
            "category": props.category,
            "content_status": props.content_status,
            "identifier": props.identifier,
            "language": props.language,
            "revision": props.revision,
            "version": props.version,
        }
        return {
            key: PptxExtractor._json_scalar(value)
            for key, value in values.items()
            if value not in {None, ""}
        }

    @staticmethod
    def _redact_metadata(metadata: dict[str, Any]) -> dict[str, Any]:
        sensitive = {"author", "last_modified_by", "comments", "identifier"}
        return {key: "[redacted]" if key in sensitive else value for key, value in metadata.items()}

    @staticmethod
    def _json_scalar(value: Any) -> float | int | str | None:
        if value is None or isinstance(value, (float, int, str)):
            return value
        if isinstance(value, (date, datetime)):
            return value.isoformat()
        return str(value)

    @staticmethod
    def _is_material_unsupported_shape(shape: Any) -> bool:
        names = {"MEDIA", "OLE_OBJECT", "LINKED_OLE_OBJECT", "WEB_VIDEO"}
        shape_type = getattr(shape, "shape_type", None)
        return getattr(shape_type, "name", "") in names
