from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def build_sample_deck(path: Path) -> Path:
    presentation = Presentation()
    presentation.core_properties.title = "Synthetic quarterly review"
    presentation.core_properties.author = "Private Author"

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly review"
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(4.0), Inches(0.8))
    paragraph = text_box.text_frame.paragraphs[0]
    paragraph.text = "Revenue increased"
    linked = text_box.text_frame.add_paragraph()
    linked.level = 1
    run = linked.add_run()
    run.text = "Evidence"
    run.hyperlink.address = "https://example.com/evidence"

    table_shape = slide.shapes.add_table(2, 2, Inches(0.7), Inches(2.2), Inches(4), Inches(1.1))
    table_shape.table.cell(0, 0).text = "Metric"
    table_shape.table.cell(0, 1).text = "Value"
    table_shape.table.cell(1, 0).text = "ARR"
    table_shape.table.cell(1, 1).text = "42"

    chart_data = ChartData()
    chart_data.categories = ["Q1", "Q2"]
    chart_data.add_series("Revenue", (30, 42))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.0),
        Inches(1.3),
        Inches(4.0),
        Inches(2.4),
        chart_data,
    )

    image_path = path.with_suffix(".png")
    Image.new("RGB", (64, 48), color=(32, 92, 160)).save(image_path)
    picture = slide.shapes.add_picture(
        str(image_path), Inches(0.7), Inches(4.0), Inches(1.6), Inches(1.2)
    )
    picture._element.nvPicPr.cNvPr.set("descr", "Blue test image")
    slide.notes_slide.notes_text_frame.text = "Confidential speaker note"

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Appendix"
    second._element.set("show", "0")
    second.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.5), Inches(1.6), Inches(1.2))
    presentation.save(path)
    image_path.unlink()
    return path
